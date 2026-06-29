#!/usr/bin/env python3
"""Konvertiert die powerpoint_*.md Decks in EINE .pptx mit nativen Tabellen,
Bullets, Monospace-Codeblöcken und echten Sprechernotizen."""
import re, glob, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Deck-Reihenfolge + Sektionsnamen (logische Vortragsordnung) ----
DECKS = [
    ("powerpoint_plan_existing_work.md",   "Related Work & Forschungslücke"),
    ("powerpoint_datasets.md",             "Datensätze in der Literatur"),
    ("powerpoint_corpora_recherche.md",    "Korpus-Auswahl"),
    ("powerpoint_phonetic_transcription.md","Transkriptionssysteme"),
    ("powerpoint_transcription_tools.md",  "Bestehende Tools (ohne LLM)"),
    ("powerpoint_evaluation_metrics.md",   "Metriken — Übersicht"),
    ("powerpoint_metric_methodology.md",   "Metrik-Methodik — Detail"),
    ("powerpoint_experiment_models.md",    "Modelle & Kosten"),
    ("powerpoint_methodik.md",             "Experiment-Methodik"),
]

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
GREY = RGBColor(0x44, 0x44, 0x44)
LIGHT = RGBColor(0xE8, 0xEE, 0xF4)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def clean_inline(t):
    t = re.sub(r'\*\*(.+?)\*\*', r'\1', t)
    t = re.sub(r'`(.+?)`', r'\1', t)
    t = re.sub(r'^\s*>\s?', '', t)
    return t.strip()

def add_title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), SW - Inches(0.8), Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = NAVY
    # Trennlinie
    ln = slide.shapes.add_shape(1, Inches(0.4), Inches(1.02), SW - Inches(0.8), Pt(2))
    ln.fill.solid(); ln.fill.fore_color.rgb = NAVY; ln.line.fill.background()
    return Inches(1.18)

def add_divider(section):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.8), Inches(3.0), SW - Inches(1.6), Inches(1.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = section
    r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)

def add_table(slide, rows, top):
    if not rows: return top
    ncol = max(len(r) for r in rows)
    rows = [r + ['']*(ncol-len(r)) for r in rows]
    nrow = len(rows)
    width = SW - Inches(0.8)
    height = Inches(min(0.34*nrow, 5.6))
    gf = slide.shapes.add_table(nrow, ncol, Inches(0.4), top, width, height)
    tbl = gf.table
    for ci in range(ncol):
        tbl.columns[ci].width = int(width/ncol)
    for ri, row in enumerate(rows):
        for ci in range(ncol):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.05); cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            para = tf.paragraphs[0]
            rn = para.add_run(); rn.text = clean_inline(row[ci])
            rn.font.size = Pt(10)
            if ri == 0:
                rn.font.bold = True; rn.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF) if ri%2 else LIGHT
    return top + height + Inches(0.15)

def add_bullets(slide, items, top):
    if not items: return top
    h = Inches(min(0.32*len(items)+0.1, 5.5))
    tb = slide.shapes.add_textbox(Inches(0.45), top, SW - Inches(0.9), h)
    tf = tb.text_frame; tf.word_wrap = True
    for i,(txt,lvl,kind) in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.level = min(lvl,4)
        r = p.add_run()
        if kind=='quote':
            r.text = '„'+txt+'"'; r.font.italic=True; r.font.size=Pt(12); r.font.color.rgb=GREY
        elif kind=='note':
            r.text = txt; r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=NAVY
        else:
            r.text = ('•  ' if lvl==0 else '–  ')+txt; r.font.size=Pt(13)
        p.space_after = Pt(3)
    return top + h + Inches(0.1)

def add_code(slide, lines, top):
    if not lines: return top
    while lines and not lines[0].strip(): lines.pop(0)
    while lines and not lines[-1].strip(): lines.pop()
    if not lines: return top
    h = Inches(min(0.165*len(lines)+0.15, 5.6))
    tb = slide.shapes.add_textbox(Inches(0.45), top, SW - Inches(0.9), h)
    tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor(0xF4,0xF6,0xF8)
    tf = tb.text_frame; tf.word_wrap = False
    for i,ln in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln.rstrip()
        r.font.name = 'Consolas'; r.font.size = Pt(9); r.font.color.rgb = RGBColor(0x22,0x33,0x44)
        p.space_after = Pt(0); p.line_spacing = 1.0
    return top + h + Inches(0.1)

# ---- Slice ein Deck in Folien ----
def parse_deck(path):
    lines = open(path, encoding='utf-8').read().splitlines()
    slides = []; cur = None
    def is_slide_head(l):
        m = re.match(r'^#{1,4}\s+(.*)$', l)
        return m and re.match(r'^Folie\b', m.group(1).strip())
    for l in lines:
        if is_slide_head(l):
            if cur: slides.append(cur)
            head = re.match(r'^#{1,4}\s+(.*)$', l).group(1).strip()
            cur = {'head': head, 'body': []}
        elif cur is not None:
            # neue Top-Sektion (# TEIL ...) beendet Folienstrom nur, wenn keine Folie folgt – wir lassen body weiterlaufen
            cur['body'].append(l)
    if cur: slides.append(cur)
    return slides

def render_slide(sd):
    s = prs.slides.add_slide(BLANK)
    head = sd['head']
    # Titel: bevorzugt **Titel:** / **Folientitel:**
    title = re.sub(r'^Folie\s+[A-Za-z0-9\-]+\s*[—:–-]\s*', '', head).strip() or head
    body = sd['body']
    # Suche expliziten Titel
    for l in body:
        m = re.match(r'\*\*(?:Titel|Folientitel):\*\*\s*(.+)', l.strip())
        if m: title = clean_inline(m.group(1)); break
    top = add_title(s, title)

    # Sprechernotiz extrahieren (alles ab Sprechernotiz/Sprecher-Notizen bis Ende/nächstes Label)
    notes = []
    in_notes = False
    content = []
    for l in body:
        st = l.strip()
        if re.match(r'(>\s*)?\*\*(Sprechernotiz|Sprecher-Notizen):\*\*', st):
            in_notes = True
            rest = re.sub(r'(>\s*)?\*\*(Sprechernotiz|Sprecher-Notizen):\*\*','',st).strip()
            if rest: notes.append(clean_inline(rest))
            continue
        if in_notes:
            if st.startswith('---') or re.match(r'^#{1,4}\s', l):
                in_notes = False
            else:
                if st: notes.append(clean_inline(st))
                continue
        content.append(l)

    # Content blockweise rendern
    i = 0; n = len(content)
    pending_bullets = []
    def flush_bullets():
        nonlocal top, pending_bullets
        if pending_bullets:
            top = add_bullets(s, pending_bullets, top); pending_bullets = []
    while i < n:
        l = content[i]; st = l.strip()
        # Label-Zeilen überspringen (aber Inhalt danach rendern)
        if re.match(r'^\*\*(Bullets|Tabelle \(Folie\)|Visuelles? .*|Tabelle.*|Stichprobe.*|Modelle.*|Metriken.*|Aufwand.*|6 Format.*|Warum.*|Erklärung.*|Beispiel.*):\*\*\s*$', st):
            i += 1; continue
        # Code-Fence
        if st.startswith('```'):
            flush_bullets(); j=i+1; code=[]
            while j<n and not content[j].strip().startswith('```'):
                code.append(content[j]); j+=1
            top = add_code(s, code, top); i=j+1; continue
        # Tabelle
        if st.startswith('|'):
            flush_bullets(); rows=[]; j=i
            while j<n and content[j].strip().startswith('|'):
                cells=[c.strip() for c in content[j].strip().strip('|').split('|')]
                if not re.match(r'^[\s:|-]+$', content[j].strip().replace('|','')):
                    rows.append(cells)
                j+=1
            top = add_table(s, rows, top); i=j; continue
        # Sichtbares Zitat / Hinweis / Kernaussage (Blockquote-Label)
        m = re.match(r'^\*\*(Sichtbares? .*|Kernaussage|Sichtbares Hinweis-Kästchen|Hypothese|Erkenntnis|Kernaussage.*):\*\*\s*(.*)', st)
        if m:
            flush_bullets()
            if m.group(2): pending_bullets.append((clean_inline(m.group(2)),0,'note'))
            i+=1; continue
        # Blockquote-Zitat
        if st.startswith('>'):
            q = clean_inline(st)
            if q: pending_bullets.append((q,0,'quote'))
            i+=1; continue
        # Bullet
        mb = re.match(r'^(\s*)[-*]\s+(.*)', l)
        if mb:
            lvl = len(mb.group(1))//2
            pending_bullets.append((clean_inline(mb.group(2)),lvl,'bullet'))
            i+=1; continue
        # sonstiger nichtleerer Text
        if st and not st.startswith('#') and not st.startswith('**'):
            pending_bullets.append((clean_inline(st),0,'bullet'))
        i+=1
    flush_bullets()

    if notes:
        s.notes_slide.notes_text_frame.text = ' '.join(notes)

# ---- Titelfolie ----
s0 = prs.slides.add_slide(BLANK)
bg = s0.shapes.add_shape(1,0,0,SW,SH); bg.fill.solid(); bg.fill.fore_color.rgb=NAVY; bg.line.fill.background()
tb = s0.shapes.add_textbox(Inches(0.8), Inches(2.4), SW-Inches(1.6), Inches(2.6)); tf=tb.text_frame; tf.word_wrap=True
p=tf.paragraphs[0]; r=p.add_run(); r.text="Assessing Multimodal Foundation Models for\nPhonetic and Tonal Transcription of Mandarin Speech"
r.font.size=Pt(28); r.font.bold=True; r.font.color.rgb=RGBColor(0xFF,0xFF,0xFF)
p2=tf.add_paragraph(); r2=p2.add_run(); r2.text="Stefan Dosch · Betreuer: Tim Schlippe · IU · Arbeitsgrundlage (auto-generiert aus Markdown)"
r2.font.size=Pt(14); r2.font.color.rgb=RGBColor(0xC8,0xD4,0xE4)

total = 0
for path, section in DECKS:
    if not os.path.exists(path): continue
    add_divider(section)
    for sd in parse_deck(path):
        render_slide(sd); total += 1

out = "Folien_aus_Markdown_AUTO.pptx"
prs.save(out)
print(f"OK: {len(prs.slides.slides._sldIdLst)} Folien -> {out}  (davon {total} Inhaltsfolien)")
